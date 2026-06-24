"""Collect every figure into a single PowerPoint deck for easy review.

One figure per slide, with a title and a short analysis caption, grouped into:
  - paper Figure 6 reproduction (the 6 MODTRAN atmospheres)
  - original RADIS-vs-reference validation (gcell / aspect)
  - all-6-atmospheres sweep
  - India / present-day comparison
plus a title slide, an overview slide and a results-summary slide.

Run after the figures exist:
    python plot_fig6_profiles.py
    python run_all_atmospheres.py
    python india_case.py
    python make_slides.py            ->  RADIS_Atmosphere_Slides.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

FIGS = "figures"
OUT = "RADIS_Atmosphere_Slides.pptx"

SW, SH = Inches(13.333), Inches(7.5)         # 16:9
DARK = RGBColor(0x1F, 0x2D, 0x3D)
GRAY = RGBColor(0x55, 0x55, 0x55)
ACCENT = RGBColor(0x1B, 0x5E, 0x9A)

# (figure file, slide title, analysis caption).  Section headers are ("#", title, subtitle).
SLIDES = [
    ("@title", "Atmosphere Models for Line-by-Line O₂ Absorption",
     "From six generic 1970s atmospheres to a real, present-day, India-specific "
     "atmosphere  —  RADIS validation of Korkin et al. (2025), JQSRT 337, 109345"),

    ("@overview", "", ""),   # special: text overview slide

    ("#", "Part 1  —  Paper Figure 6:  the six MODTRAN atmospheres",
     "All 6 standard models extracted from the reference code's hprofiles.h "
     "(pressure, temperature, density + 8 gases at 50 heights)."),
    ("fig6a_height_grid.png", "Fig. 6(a)  —  the MODTRAN height grid",
     "50 levels in three bands: 1 km steps below 25 km, 2.5 km to 50 km, 5 km to "
     "120 km. Finest where most of the atmosphere (and absorption) sits."),
    ("fig6b_PTD_profiles.png", "Fig. 6(b)  —  pressure, temperature, density",
     "Temperature (middle) varies strongly between the six models; pressure and "
     "density (sides) barely differ — exactly as the paper states."),
    ("fig6c_H2O_CO2.png", "Fig. 6(c)  —  water vapour & CO₂",
     "Water vapour spans orders of magnitude between tropical and sub-arctic and "
     "drops fast with height; CO₂ is well-mixed."),
    ("fig6d_O3_N2O.png", "Fig. 6(d)  —  ozone & N₂O",
     "Ozone peaks in the stratosphere (~25-30 km); N₂O is highest near the "
     "ground and falls off with altitude."),
    ("fig6e_CO_CH4.png", "Fig. 6(e)  —  CO & methane",
     "CO rises again high up; methane is well-mixed low down and decays in the "
     "stratosphere."),
    ("fig6f_O2_NO2.png", "Fig. 6(f)  —  oxygen & NO₂",
     "Oxygen is well-mixed at 20.9% (nearly identical across all models) — "
     "this is what makes the O₂ A-band a clean reference."),

    ("#", "Part 2  —  Original RADIS-vs-reference validation",
     "Our line-by-line code (RADIS) overplotted on the paper's own C-code outputs."),
    ("overplot_o2a_gcell_vs_radis.png", "Gas cell: O₂ A-band",
     "RADIS (dashed) vs the paper's gcell C-code (solid) for the O₂ A-band: "
     "match within ~2.4%."),
    ("overplot_ch4_gcell_vs_radis.png", "Gas cell: methane 2.3 µm",
     "Same check for the CH₄ 2.3 µm band: agreement within ~8%, the "
     "residual explained by HITRAN version and line-wing cutoff."),
    ("overplot_aspect_o2a_radis_stack.png", "Atmosphere: RADIS slab-stack vs aspect",
     "Partial-column optical depth (top-of-atmosphere down to 0/1/2.5/8 km). RADIS "
     "stack (dashed) reproduces the paper's aspect C-code (solid) within 2.7-4.5%."),
    ("overplot_aspect_o2a_fullcolumn_zoom.png", "Full-column detail + residual",
     "Zoom on the band head with the RADIS-minus-reference residual: small and "
     "structureless, i.e. no modelling disagreement."),
    ("aspect_o2a_height_variation.png", "Absorption vs height",
     "Band-integrated O₂ absorption falls smoothly as the observation point "
     "rises — less air (and O₂) remains above."),

    ("#", "Part 3  —  Running all six standard atmospheres",
     "The same general code path now handles every MODTRAN model, not just US-1976."),
    ("all6_O2_o2a_optical_depth.png", "O₂ A-band: all six models overlaid",
     "Full-column optical depth for all six atmospheres. Curves nearly coincide "
     "because oxygen is well-mixed; the point is one code runs every case."),
    ("all6_O2_o2a_bandintegral.png", "Total absorption across the six models",
     "Band-integrated absorption per model: they agree within ~1%, confirming the "
     "O₂ column barely changes between climates."),

    ("#", "Part 4  —  India / present-day atmosphere (points 2 & 3)",
     "Replace the generic model with the REAL atmosphere over Mumbai (IIT Bombay), "
     "from NASA PSG (GEOS-5) and NRLMSIS 2.1."),
    ("india_profiles_Mumbai_IIT_Bombay.png", "Mumbai: generic vs real atmosphere",
     "Present-day profiles (PSG = blue, NRLMSIS = green) agree with each other and "
     "differ from the 1970s 'Tropical' model (grey), notably near the 17 km "
     "tropopause and in the stratosphere."),
    ("india_o2a_optical_depth_Mumbai_IIT_Bombay.png",
     "O₂ A-band over Mumbai: generic vs present-day",
     "The real atmosphere absorbs slightly LESS than the generic model predicts. "
     "Quantified on the summary slide."),

    ("@results", "", ""),   # special: results table slide
]

OVERVIEW = [
    ("What the professor asked", [
        "1.  The reference code uses MODTRAN's six standard atmospheres "
        "(Tropical, Mid-Lat Summer/Winter, Sub-Arctic Summer/Winter, US-1976). "
        "Do we have the data, and can we reproduce Figure 6?",
        "2.  Those models are generic and decades old; the real atmosphere differs. "
        "Add an exact present-day model (e.g. NASA PSG) and include the variation.",
        "3.  None of the models is India-specific. Use a model suitable for the "
        "Indian subcontinent, ideally a recent one.",
    ]),
    ("What we did", [
        "1.  Extracted all six MODTRAN atmospheres and reproduced Figure 6(a-f); "
        "the code now runs for all six, not just US-1976.",
        "2.  Added real, date- and location-specific atmospheres from NASA PSG "
        "(GEOS-5) and NRLMSIS 2.1 (2021) behind the same calculation.",
        "3.  Worked India example (Mumbai / IIT Bombay) comparing the generic "
        "Tropical model against the two present-day models for the O₂ A-band.",
    ]),
]

RESULTS_ROWS = [
    ("Atmosphere", "O₂ column (cm⁻²)", "Total absorption", "vs generic"),
    ("MODTRAN Tropical (1970s)", "4.523 ×10²⁴", "1023", "—"),
    ("NASA PSG (exact, 15 Jun 2024)", "4.408 ×10²⁴", "996", "−2.6 %"),
    ("NRLMSIS 2.1 (latest)", "4.471 ×10²⁴", "1011", "−1.2 %"),
]


def _title_box(slide, text, top=0.35, size=26, color=DARK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(top), SW - Inches(1.2),
                                  Inches(1.0))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = color
    return tb


def add_section(prs, title, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = ACCENT
    tb = s.shapes.add_textbox(Inches(0.8), Inches(2.6), SW - Inches(1.6), Inches(2))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = title
    r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = RGBColor(255, 255, 255)
    if subtitle:
        p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = subtitle
        r2.font.size = Pt(15); r2.font.color.rgb = RGBColor(0xDD, 0xE6, 0xF0)


def add_title_slide(prs, title, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = DARK
    tb = s.shapes.add_textbox(Inches(0.9), Inches(2.3), SW - Inches(1.8), Inches(3))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = title
    r.font.size = Pt(34); r.font.bold = True; r.font.color.rgb = RGBColor(255, 255, 255)
    p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = subtitle
    r2.font.size = Pt(16); r2.font.color.rgb = RGBColor(0xBC, 0xCD, 0xDE)


def add_figure(prs, fname, title, caption):
    path = os.path.join(FIGS, fname)
    if not os.path.exists(path):
        print("  (missing, skipped)", fname); return
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _title_box(s, title, top=0.3, size=24, color=ACCENT)
    pic = s.shapes.add_picture(path, 0, 0, height=Inches(5.25))
    pic.left = int((SW - pic.width) / 2)
    pic.top = Inches(1.25)
    cap = s.shapes.add_textbox(Inches(0.7), Inches(6.6), SW - Inches(1.4), Inches(0.8))
    tf = cap.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = caption
    r.font.size = Pt(13); r.font.italic = True; r.font.color.rgb = GRAY


def add_overview(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _title_box(s, "Overview", top=0.3, size=26)
    top = 1.4
    for head, items in OVERVIEW:
        tb = s.shapes.add_textbox(Inches(0.7), Inches(top), SW - Inches(1.4), Inches(2.8))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; r = p.add_run(); r.text = head
        r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = ACCENT
        for it in items:
            pp = tf.add_paragraph(); rr = pp.add_run(); rr.text = it
            rr.font.size = Pt(13); rr.font.color.rgb = DARK
            pp.space_after = Pt(4)
        top += 2.9


def add_results(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _title_box(s, "Result  —  Mumbai, 15 June 2024 (O₂ A-band)", top=0.3, size=24)
    rows, cols = len(RESULTS_ROWS), 4
    tbl = s.shapes.add_table(rows, cols, Inches(0.8), Inches(1.5),
                             SW - Inches(1.6), Inches(2.4)).table
    for c in range(cols):
        tbl.columns[c].width = Inches((SW.inches - 1.6) / cols)
    for ri, row in enumerate(RESULTS_ROWS):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci); cell.text = val
            par = cell.text_frame.paragraphs[0]
            par.runs[0].font.size = Pt(14)
            if ri == 0:
                par.runs[0].font.bold = True
                par.runs[0].font.color.rgb = RGBColor(255, 255, 255)
                cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT
    notes = ("Takeaway:  the generic 1970s 'Tropical' model OVER-predicts real "
             "O₂ A-band absorption over Mumbai by ~1-3% on this day, and the two "
             "independent present-day models (PSG, NRLMSIS) agree to ~1.5% — so the "
             "difference is real, not a quirk of one data source.")
    tb = s.shapes.add_textbox(Inches(0.8), Inches(4.6), SW - Inches(1.6), Inches(2))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = notes
    r.font.size = Pt(15); r.font.color.rgb = DARK


def main():
    prs = Presentation()
    prs.slide_width = SW; prs.slide_height = SH
    for fname, title, caption in SLIDES:
        if fname == "@title":
            add_title_slide(prs, title, caption)
        elif fname == "#":
            add_section(prs, title, caption)
        elif fname == "@overview":
            add_overview(prs)
        elif fname == "@results":
            add_results(prs)
        else:
            add_figure(prs, fname, title, caption)
    prs.save(OUT)
    print(f"Wrote {OUT}  ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
